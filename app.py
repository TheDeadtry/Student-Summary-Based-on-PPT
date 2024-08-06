from flask import Flask, render_template, request, jsonify
from pptx import Presentation
from sklearn.metrics.pairwise import cosine_similarity
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
import gensim.downloader as api
import numpy as np
import pytesseract
from PIL import Image
import io
import logging


print("Starting to import libraries...")

app = Flask(__name__)

print("Import Libraries Complete")
print("Starting nltk download")

nltk.download('punkt')
nltk.download('stopwords')

print("NLTK downloads complete")

print("Starting to load word vectors...")
# Load pre-trained word embeddings
word_vectors = api.load('glove-twitter-25')
print("Loading word vectors complete")

print("Starting function definitions")

# Set up logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# Set Tesseract path if necessary
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + " "
    return text

def preprocess_text(text):
    stop_words = set(stopwords.words('english'))
    tokens = word_tokenize(text.lower())
    return [token for token in tokens if token.isalnum() and token not in stop_words]

def text_to_vec(text, word_vectors):
    words = preprocess_text(text)
    word_vecs = [word_vectors[word] for word in words if word in word_vectors]
    if not word_vecs:
        return np.zeros(word_vectors.vector_size)
    return np.mean(word_vecs, axis=0)

def semantic_chunking(text, chunk_size=50):
    words = preprocess_text(text)
    return [' '.join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def calculate_similarity(ppt_text, summary):
    ppt_chunks = semantic_chunking(ppt_text)
    summary_vec = text_to_vec(summary, word_vectors)

    chunk_scores = []
    for chunk in ppt_chunks:
        chunk_vec = text_to_vec(chunk, word_vectors)
        similarity = cosine_similarity([chunk_vec], [summary_vec])[0][0]
        chunk_scores.append(similarity)
    print("Similarity calculation complete")

    # Take the average of the top 3 chunk scores
    top_scores = sorted(chunk_scores, reverse=True)[:3]
    avg_top_score = np.mean(top_scores) if top_scores else 0

    # Calculate the overall relevance
    relevance_threshold = 0.3  # Adjust this value as needed
    relevant_chunks = sum(1 for score in chunk_scores if score > relevance_threshold)
    relevance_ratio = relevant_chunks / len(chunk_scores) if chunk_scores else 0

    # Combine similarity and relevance
    final_score = (avg_top_score * 0.7 + relevance_ratio * 0.3) * 100

    return final_score

def extract_text_from_image(file):
    try:
        image = Image.open(io.BytesIO(file.read()))
        logger.debug(f"Image format: {image.format}, Size: {image.size}, Mode: {image.mode}")
        text = pytesseract.image_to_string(image)
        logger.debug(f"Extracted text: {text}")
        return text
    except Exception as e:
        logger.error(f"Error in extract_text_from_image: {str(e)}")
        raise

print("Function definitions complete")

@app.route('/', methods=['GET', 'POST'])
def index():
    print("Index route accessed")
    return render_template('index.html')

@app.route('/score', methods=['POST'])
def score():
    print("Score route accessed")
    if 'ppt' not in request.files:
        return jsonify({'error': 'No PowerPoint file in the request'}), 400

    ppt_file = request.files['ppt']

    if ppt_file.filename == '':
        return jsonify({'error': 'No PowerPoint file selected'}), 400

    if not ppt_file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'PowerPoint file must be a .pptx file'}), 400

    summary = ""
    if 'summary' in request.form:
        summary = request.form['summary']
    elif 'summary_image' in request.files:
        summary_image = request.files['summary_image']
        if summary_image.filename == '':
            return jsonify({'error': 'No summary image selected'}), 400
        if not summary_image.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
            return jsonify({'error': 'Summary image must be a .png, .jpg, or .jpeg file'}), 400
        summary = extract_text_from_image(summary_image)
    else:
        return jsonify({'error': 'No summary text or image provided'}), 400

    try:
        ppt_text = extract_text_from_pptx(ppt_file)
        score = calculate_similarity(ppt_text, summary)
        print(f"Calculated score: {score}")
        return jsonify({'score': round(score, 2)})
    except Exception as e:
        logger.error(f"Error occurred: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/extract_text', methods=['POST'])
def extract_text():
    if 'summary_image' not in request.files:
        return jsonify({'error': 'No image file in the request'}), 400

    summary_image = request.files['summary_image']

    if summary_image.filename == '':
        return jsonify({'error': 'No image file selected'}), 400

    if not summary_image.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
        return jsonify({'error': 'File must be a .png, .jpg, or .jpeg file'}), 400

    try:
        text = extract_text_from_image(summary_image)
        if not text.strip():
            return jsonify({'error': 'No text could be extracted from the image'}), 400
        return jsonify({'text': text})
    except Exception as e:
        logger.error(f"Error in extract_text route: {str(e)}")
        return jsonify({'error': f'Failed to process image: {str(e)}'}), 500

if __name__ == '__main__':
    print("Starting Flask app...")
    app.run(debug=True, host='0.0.0.0', port=5000)